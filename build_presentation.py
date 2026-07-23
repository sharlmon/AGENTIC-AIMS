import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Stealth Dark Navy & Vibrant Indigo
    COLOR_BG = RGBColor(15, 17, 26)         # Deep Navy #0F111A
    COLOR_CARD = RGBColor(27, 31, 48)       # Slate Card #1B1F30
    COLOR_ACCENT = RGBColor(99, 102, 241)   # Vibrant Indigo #6366F1
    COLOR_TEXT_WHITE = RGBColor(248, 250, 252) # Crisp White
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184) # Muted Gray-Blue
    COLOR_HIGHLIGHT = RGBColor(238, 242, 255) # Light Text Accent

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: Title Slide (Hook + Visual Hint)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_BG)

    # Accent decorative background shape
    shape = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(6), Inches(6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(49, 46, 129)
    shape.line.fill.background()

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(10), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "48-HOUR AI HACKATHON 2026"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "JITUME AGENCY OS"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.name = "Georgia"

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(9.5), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI that turns raw client meetings into signed-ready proposals before your coffee even cools down."
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_TEXT_MUTED

    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.2), Inches(3.8), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(30, 27, 75)
    badge.line.color.rgb = COLOR_ACCENT
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "BUILT BY TEAM AGENTIC AIMS"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_HIGHLIGHT

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(6.2), Inches(8), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Presented by Sharlmon Junior · Technical Lead"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 2: Team Slide (Name, Role & WHAT THEY BUILT)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_BG)

    txBox = slide2.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(10), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "THE TEAM"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    txBox = slide2.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(10), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Agentic AIMS — Proof of Build"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.name = "Georgia"

    # Exact proof of work for each member
    members = [
        ("Sharlmon Junior", "Technical Lead", "Built Dual-Node Pipeline & HITL Engine", "S"),
        ("Suleiman", "Lead Developer", "Built Prisma ORM & Webhook Ingestion", "S"),
        ("Joshua", "AI Engineer", "Built Gemini 1.5 Pro & NIM Fallback", "J"),
        ("Kofa", "Product Designer", "Built Zinc 950 Stealth UI Workspace", "K"),
        ("Sharlmon", "Pitch & Ops Lead", "Built Resend Email & Calendar Engine", "S"),
    ]

    card_width = Inches(2.1)
    card_gap = Inches(0.2)
    start_x = Inches(1.0)
    card_y = Inches(2.3)
    card_height = Inches(4.5)

    for idx, (name, role, built, letter) in enumerate(members):
        cx = start_x + idx * (card_width + card_gap)
        
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = RGBColor(45, 52, 78)

        avatar = slide2.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.55), card_y + Inches(0.4), Inches(1.0), Inches(1.0))
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = COLOR_ACCENT if idx == 0 else RGBColor(16, 185, 129) if idx == 1 else RGBColor(245, 158, 11)
        avatar.line.fill.background()
        
        p = avatar.text_frame.paragraphs[0]
        p.text = letter
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_WHITE

        tx = slide2.shapes.add_textbox(cx + Inches(0.05), card_y + Inches(1.5), card_width - Inches(0.1), Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_WHITE

        tx = slide2.shapes.add_textbox(cx + Inches(0.05), card_y + Inches(2.1), card_width - Inches(0.1), Inches(0.5))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = role
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT

        tx = slide2.shapes.add_textbox(cx + Inches(0.1), card_y + Inches(2.7), card_width - Inches(0.2), Inches(1.6))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Built: {built}"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 3: Problem Statement (Specific Scenario & Numbers)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_BG)

    txBox = slide3.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(10), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "THE PROBLEM"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    txBox = slide3.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11), Inches(1.0))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Onboarding a client shouldn't feel like a second job."
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.name = "Georgia"

    card1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.4), Inches(4.5), Inches(4.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD
    card1.line.color.rgb = RGBColor(239, 68, 68)

    txBox = slide3.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(4.1), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "10+ HRS"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(251, 191, 36)

    txBox = slide3.shapes.add_textbox(Inches(1.2), Inches(4.7), Inches(4.1), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Agencies lose 10+ hours per week on manual meeting notes, calendar scheduling, and proposal drafting."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_MUTED

    card2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.0), Inches(2.4), Inches(6.3), Inches(4.3))
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_CARD
    card2.line.color.rgb = RGBColor(45, 52, 78)

    pain_points = [
        ("1", "Manual Call Transcribing: Hours wasted typing summary notes by hand"),
        ("2", "Scheduling Friction: A slow back-and-forth fight with Google Calendar"),
        ("3", "Communication Delays: Slow follow-up emails lead to lost client interest"),
        ("4", "Manual Proposal Creation: 4-hour drafting of proposals & invoice tables"),
    ]

    for idx, (num, text) in enumerate(pain_points):
        py = Inches(2.7) + idx * Inches(0.95)
        num_circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.4), py, Inches(0.45), Inches(0.45))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLOR_ACCENT
        num_circle.line.fill.background()
        p = num_circle.text_frame.paragraphs[0]
        p.text = num
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_WHITE

        tx = slide3.shapes.add_textbox(Inches(7.0), py - Inches(0.05), Inches(5.0), Inches(0.6))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_WHITE

    # ==========================================
    # SLIDE 4: Solution (One Clear Sentence)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_BG)

    txBox = slide4.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(10), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "THE SOLUTION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    txBox = slide4.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Jitume Agency OS is an autonomous dual-node AI system that turns raw discovery calls into signed-ready proposals in minutes."
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.name = "Georgia"

    # Solution Banner Box
    card_sol = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.2), Inches(11.3), Inches(3.5))
    card_sol.fill.solid()
    card_sol.fill.fore_color.rgb = COLOR_CARD
    card_sol.line.color.rgb = COLOR_ACCENT

    tx = slide4.shapes.add_textbox(Inches(1.4), Inches(3.6), Inches(10.5), Inches(2.6))
    tf = tx.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "⚡ Phase 1: Zero-Touch Automation — Ingests calls, schedules team syncs, and emails clients automatically."
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(16, 185, 129)

    p2 = tf.add_paragraph()
    p2.text = "\n🛡️ Phase 2: Human-in-the-Loop Control — Gemini AI synthesizes proposals & HTML invoices with 1-click admin approval."
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_HIGHLIGHT

    # ==========================================
    # SLIDE 5: MVP / Core Features (3 User Capabilities Max)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_BG)

    txBox = slide5.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(10), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "MVP FEATURES"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    txBox = slide5.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = "3 Core Capabilities What Users Can Do"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.name = "Georgia"

    features = [
        ("1. Zero-Touch Discovery & Scheduling", "Client calls are automatically transcribed into Contact Reports, and a 10:00 AM team sync with Google Meet is booked instantly."),
        ("2. AI Proposal & Invoice Synthesis", "Multi-transcript synthesis generates complete proposals and formatted HTML invoice tables using Gemini 1.5 Pro."),
        ("3. Mission Control HITL Review Queue", "Admins refine call notes inline, edit HTML markup live, and approve final email dispatches with a single click."),
    ]

    for idx, (title, desc) in enumerate(features):
        fx = Inches(1.0) + idx * Inches(3.9)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fx, Inches(2.4), Inches(3.6), Inches(4.3))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = RGBColor(45, 52, 78)

        tx = slide5.shapes.add_textbox(fx + Inches(0.2), Inches(2.7), Inches(3.2), Inches(0.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT

        tx = slide5.shapes.add_textbox(fx + Inches(0.2), Inches(3.7), Inches(3.2), Inches(2.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_WHITE

    # ==========================================
    # SLIDE 6: Live Demo (Win or Lose the Room)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_BG)

    txBox = slide6.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(10), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "LIVE DEMONSTRATION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)

    txBox = slide6.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Live Onboarding Walkthrough"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.name = "Georgia"

    card_demo = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.3), Inches(11.3), Inches(4.4))
    card_demo.fill.solid()
    card_demo.fill.fore_color.rgb = COLOR_CARD
    card_demo.line.color.rgb = RGBColor(16, 185, 129)

    steps = [
        "Step 1: Init Project ('george' / 'E-commerce Platform Development')",
        "Step 2: Fathom Webhook Ingestion & Zero-Touch Contact Report",
        "Step 3: Google Calendar Auto-Scheduler (Tomorrow 10:00 AM Sync)",
        "Step 4: Mission Control HITL Review & One-Click Resend Dispatch",
    ]

    for idx, step in enumerate(steps):
        tx = slide6.shapes.add_textbox(Inches(1.4), Inches(2.7) + idx * Inches(0.9), Inches(10.5), Inches(0.7))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = f"▶ {step}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_WHITE

    # ==========================================
    # SLIDE 7: Tech Stack & Architecture
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLOR_BG)

    txBox = slide7.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(10), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "SYSTEM ARCHITECTURE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    txBox = slide7.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Production-Grade Tech Stack"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.name = "Georgia"

    stacks = [
        ("Frontend UI", "Next.js 14 App Router, React 18, TypeScript, Tailwind v4 Stealth Theme"),
        ("AI LLM Engines", "Primary: Google Gemini 1.5 Pro\nFallback: NVIDIA NIM (Llama 3.3 70B)"),
        ("Persistence", "Prisma ORM v6 with embedded SQLite database (dev.db)"),
        ("Integrations", "Fathom Webhooks, Google Calendar API (Service Auth), Resend API"),
    ]

    for idx, (title, desc) in enumerate(stacks):
        rx = Inches(1.0) if idx % 2 == 0 else Inches(6.9)
        ry = Inches(2.3) if idx < 2 else Inches(4.6)
        
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, ry, Inches(5.4), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = RGBColor(45, 52, 78)

        tx = slide7.shapes.add_textbox(rx + Inches(0.3), ry + Inches(0.2), Inches(4.8), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT

        tx = slide7.shapes.add_textbox(rx + Inches(0.3), ry + Inches(0.7), Inches(4.8), Inches(1.1))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_WHITE

    # ==========================================
    # SLIDE 8: Risk & Impact Strategy
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLOR_BG)

    txBox = slide8.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(10), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "IMPACT & RISK STRATEGY"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    txBox = slide8.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11), Inches(0.8))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Addressing Risks & Business Impact"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.name = "Georgia"

    card_r1 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.3), Inches(5.4), Inches(4.4))
    card_r1.fill.solid()
    card_r1.fill.fore_color.rgb = COLOR_CARD
    card_r1.line.color.rgb = RGBColor(239, 68, 68)

    tx = slide8.shapes.add_textbox(Inches(1.3), Inches(2.6), Inches(4.8), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    p.text = "BIGGEST RISKS IDENTIFIED"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)

    r_bullets = [
        "LLM Hallucinations in Pricing/Scope",
        "API Rate Limits during Peak Loads",
        "Client Email Sandbox Restrictions",
    ]
    for idx, bullet in enumerate(r_bullets):
        tx = slide8.shapes.add_textbox(Inches(1.3), Inches(3.3) + idx * Inches(1.0), Inches(4.8), Inches(0.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"⚠️ {bullet}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_WHITE

    card_r2 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.3), Inches(5.4), Inches(4.4))
    card_r2.fill.solid()
    card_r2.fill.fore_color.rgb = COLOR_CARD
    card_r2.line.color.rgb = RGBColor(16, 185, 129)

    tx = slide8.shapes.add_textbox(Inches(7.2), Inches(2.6), Inches(4.8), Inches(0.5))
    p = tx.text_frame.paragraphs[0]
    p.text = "MITIGATION & PROOF"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(16, 185, 129)

    m_bullets = [
        "HITL Review Gate ensures 100% human oversight",
        "NVIDIA NIM (Llama 3.3 70B) auto-failover engine",
        "Verified domain fallback (Sharlmon <hello@sharl-tech.co.ke>)",
    ]
    for idx, bullet in enumerate(m_bullets):
        tx = slide8.shapes.add_textbox(Inches(7.2), Inches(3.3) + idx * Inches(1.0), Inches(4.8), Inches(0.8))
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"🛡️ {bullet}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_TEXT_WHITE

    # ==========================================
    # SLIDE 9: Conclusion / Q&A Slide
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, COLOR_BG)

    txBox = slide9.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "JITUME AGENCY OS"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.font.name = "Georgia"

    txBox = slide9.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Turning client meetings into signed contracts in minutes."
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.color.rgb = COLOR_TEXT_MUTED

    badge = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.66), Inches(4.8), Inches(4.0), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_ACCENT
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU — READY FOR Q&A"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE

    output_filename = "Jitume_Agency_OS_Pitch_v3.pptx"
    prs.save(output_filename)
    print(f"Successfully generated official 9-slide deck: {output_filename}")

if __name__ == "__main__":
    create_presentation()
